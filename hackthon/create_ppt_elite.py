from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os

# ── Paths to generated images ────────────────────────────
BRAIN = r"C:\Users\vatti\.gemini\antigravity\brain\e3d243a0-ea5e-4497-822c-217f567afa1c"
IMG_COVER = os.path.join(BRAIN, "janrakshak_cover_professional_1774271466444.png")
IMG_APP   = os.path.join(BRAIN, "janrakshak_solution_mockup_1774271481988.png")
IMG_TECH  = os.path.join(BRAIN, "janrakshak_tech_architecture_v3_1774271537786.png")

# ── Colors ───────────────────────────────────────────────
C_PRIMARY    = RGBColor(0x06,0x4e,0x3b) # Deep Emerald
C_SECONDARY  = RGBColor(0x10,0xb9,0x81) # Emerald
C_ACCENT     = RGBColor(0x34,0xd3,0x99) # Lite Emerald
C_HIGHLIGHT  = RGBColor(0xfb,0xbf,0x24) # Amber
C_DANGER     = RGBColor(0xef,0x44,0x44) # Red
C_NAVY       = RGBColor(0x0f,0x17,0x2a) # Dark Navy
C_WHITE      = RGBColor(0xff,0xff,0xff)
C_GRAY       = RGBColor(0x94,0xa3,0xb8)

W, H = Inches(13.33), Inches(7.5)

prs = Presentation()
prs.slide_width, prs.slide_height = W, H

def blank_slide(): return prs.slides.add_slide(prs.slide_layouts[6])

def add_rect(slide, l,t,w,h, fill=None, line=None):
    s = slide.shapes.add_shape(1, l, t, w, h)
    s.line.width = Pt(0)
    if fill: s.fill.solid(); s.fill.fore_color.rgb = fill
    else: s.fill.background()
    return s

def add_txt(slide, txt, l,t,w,h, sz=18, bold=False, col=C_WHITE, align=PP_ALIGN.LEFT, italic=False):
    tb = slide.shapes.add_textbox(l,t,w,h)
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = str(txt)
    r.font.size = Pt(sz); r.font.bold = bold; r.font.italic = italic
    r.font.color.rgb = col
    return tb

def add_img(slide, path, l,t,w,h):
    if os.path.exists(path): slide.shapes.add_picture(path, l,t,w,h)

# ──────────────────────────────────────────────────────────
# SLIDE 1: ELITE TITLE
# ──────────────────────────────────────────────────────────
s = blank_slide()
add_img(s, IMG_COVER, 0, 0, W, H)
add_rect(s, 0, 0, W*0.55, H, fill=RGBColor(0x06,0x1a,0x12))
add_rect(s, 0, 0, Inches(0.12), H, fill=C_SECONDARY)

add_txt(s, "🏆 National Level Hackathon 2026", Inches(0.6), Inches(0.8), Inches(6), Inches(0.5), sz=14, bold=True, col=C_HIGHLIGHT)
add_txt(s, "🌱 JanRakshak AI", Inches(0.6), Inches(1.8), Inches(7), Inches(1.5), sz=64, bold=True)
add_txt(s, '"Predict Before It\'s Critical"', Inches(0.6), Inches(3.2), Inches(6.5), Inches(0.6), sz=24, col=C_ACCENT, italic=True)
add_txt(s, "Hyper-Resilient, Offline-First Medical Intelligence for\nIndia's 1.2 Million ASHA Workers & 600M Rural Citizens", 
        Inches(0.6), Inches(4.2), Inches(6), Inches(1), sz=18, col=C_WHITE)

stats = [("600M+","Rural Target"), ("<100ms","AI Latency"), ("Zero","Cloud Dependency"), ("3","Indian Langs")]
for i, (n, l) in enumerate(stats):
    xl = Inches(0.6 + i*1.6)
    add_rect(s, xl, Inches(5.8), Inches(1.4), Inches(1.2), fill=C_PRIMARY)
    add_txt(s, n, xl, Inches(5.95), Inches(1.4), Inches(0.4), sz=20, bold=True, col=C_ACCENT, align=PP_ALIGN.CENTER)
    add_txt(s, l, xl, Inches(6.5), Inches(1.4), Inches(0.3), sz=10, col=C_WHITE, align=PP_ALIGN.CENTER)

# ──────────────────────────────────────────────────────────
# SLIDE 2: THE PROBLEM (Data-Driven Pan-India Crisis)
# ──────────────────────────────────────────────────────────
s = blank_slide()
add_rect(s, 0, 0, W, H, fill=C_NAVY)
add_rect(s, 0, 0, W, Inches(0.12), fill=C_DANGER)

add_txt(s, "🔴 The Problem: India's Rural Health Blackhole", Inches(0.5), Inches(0.4), Inches(12), Inches(0.8), sz=38, bold=True)
add_txt(s, "600 million Indians live in the shadow of 'Late Detection'—where illness only meeting reality after it's too late.", 
        Inches(0.5), Inches(1.2), Inches(12), Inches(0.4), sz=16, col=C_ACCENT, italic=True)

problems = [
    ("1.4 Million", "Preventable deaths annually\ndue to delayed detection", C_DANGER),
    ("95% Access Gap", "ASHA workers still usage paper\nregisters for critical vitals", C_HIGHLIGHT),
    ("30+ KM Journey", "Average distance to nearest\nER-capable hospital", C_SECONDARY),
    ("Zero Signal", "Villages with no stable 4G/5G\nremain digitally invisible", C_ACCENT)
]

for i, (n, l, c) in enumerate(problems):
    xl = Inches(0.5 + i*3.1)
    add_rect(s, xl, Inches(2.2), Inches(2.9), Inches(2.5), fill=RGBColor(0x1e,0x29,0x3b))
    add_rect(s, xl, Inches(2.2), Inches(2.9), Inches(0.1), fill=c)
    add_txt(s, n, xl, Inches(2.6), Inches(2.9), Inches(0.8), sz=30, bold=True, col=c, align=PP_ALIGN.CENTER)
    add_txt(s, l, xl, Inches(3.6), Inches(2.9), Inches(0.6), sz=12, col=C_WHITE, align=PP_ALIGN.CENTER)

add_rect(s, Inches(0.5), Inches(5.2), Inches(12.3), Inches(1.5), fill=RGBColor(0x2d,0x10,0x10))
add_txt(s, "Current Workflow Error: Paper Register → No Risk Scoring → Delayed ER Routing → Fatality", 
        Inches(1.0), Inches(5.4), Inches(11.3), Inches(0.5), sz=18, bold=True, col=C_DANGER, align=PP_ALIGN.CENTER)
add_txt(s, "Goal: Bypass the broken system and bring clinical-grade AI to a ₹6000 smartphone.", 
        Inches(1.0), Inches(6.1), Inches(11.3), Inches(0.4), sz=14, col=C_GRAY, align=PP_ALIGN.CENTER)

# ──────────────────────────────────────────────────────────
# SLIDE 3: OUR SOLUTION (Elite Feature Set)
# ──────────────────────────────────────────────────────────
s = blank_slide()
add_rect(s, 0, 0, W, H, fill=C_PRIMARY)
add_img(s, IMG_APP, Inches(7.5), Inches(0.8), Inches(5.5), Inches(5.8))

add_txt(s, "💡 Solution: JanRakshak AI — 'Pocket PHC'", Inches(0.5), Inches(0.4), Inches(7), Inches(0.8), sz=38, bold=True)
add_txt(s, "Decentralizing healthcare through localized AI intelligence.", Inches(0.5), Inches(1.2), Inches(7), Inches(0.4), sz=16, col=C_ACCENT, italic=True)

feats = [
    ("🧠 AI Risk Scoring", "Heart, Diabetes, and Anemia metrics analyzed under 100ms."),
    ("🤱 Maternal Care", "Automated High-Risk Pregnancy identification & protocol."),
    ("📴 Hyper-Resilience", "Works 100% Offline; auto-syncs or sends Alerts via SMS."),
    ("🗺️ GPS Smart Route", "Dynamic hospital mapping based on current patient location."),
    ("💬 Voice & Lang", "English, Hindi, Telugu voice feedback for illiterate patients."),
    ("💊 Medicine Engine", "Fuzzy-match search for 20+ diseases & 80+ home remedies.")
]

for i, (t, b) in enumerate(feats):
    row, col = i // 2, i % 2
    yl = Inches(2.2 + row*1.6)
    xl = Inches(0.5 + col*3.4)
    add_rect(s, xl, yl, Inches(3.2), Inches(1.4), fill=RGBColor(0x06,0x5f,0x46))
    add_txt(s, t, xl+Inches(0.1), yl+Inches(0.1), Inches(3), Inches(0.4), sz=14, bold=True, col=C_ACCENT)
    add_txt(s, b, xl+Inches(0.1), yl+Inches(0.5), Inches(3), Inches(0.7), sz=11, col=C_WHITE)

# ──────────────────────────────────────────────────────────
# SLIDE 4: TECHNICAL DEPTH (Engine Room)
# ──────────────────────────────────────────────────────────
s = blank_slide()
add_rect(s, 0, 0, W, H, fill=C_NAVY)
add_img(s, IMG_TECH, Inches(0.5), Inches(1.0), Inches(6.5), Inches(5.5))

add_txt(s, "🏗️ Technical Depth: High-Performance PWA Architecture", Inches(0.5), Inches(0.4), Inches(12), Inches(0.8), sz=34, bold=True)

tech_grid = [
    ("Client", "React 18 + Vite", "Zero-lag UI on ₹6000 Android phones."),
    ("Server", "Python Flask API", "Lightweight, secure, and easily scalable."),
    ("Vision", "Gemini 1.5 Flash", "Multimodal wound & fracture diagnosis."),
    ("GIS", "Leaflet.js + OSM", "Offline routing using current GPS anchors."),
    ("i18n", "Custom Engine", "0 libraries; instant Lang-switching (Hi/Te/En)."),
    ("Comm", "tel: protocol", "Direct-to-dial emergency bypass logic.")
]

for i, (l, t, n) in enumerate(tech_grid):
    yl = Inches(1.2 + i*1.0)
    add_rect(s, Inches(7.3), yl, Inches(5.5), Inches(0.9), fill=RGBColor(0x1e,0x29,0x3b))
    add_rect(s, Inches(7.3), yl, Inches(0.1), Inches(0.9), fill=C_ACCENT)
    add_txt(s, l, Inches(7.5), yl+Inches(0.1), Inches(1.5), Inches(0.4), sz=10, col=C_ACCENT, bold=True)
    add_txt(s, t, Inches(7.5), yl+Inches(0.35), Inches(3.5), Inches(0.4), sz=14, bold=True)
    add_txt(s, n, Inches(10.0), yl+Inches(0.2), Inches(2.8), Inches(0.6), sz=11, col=C_GRAY)

# ──────────────────────────────────────────────────────────
# SLIDE 5: INNOVATION & UNIQUENESS
# ──────────────────────────────────────────────────────────
s = blank_slide()
add_rect(s, 0, 0, W, H, fill=C_PRIMARY)
add_txt(s, "🚀 What Makes Us Truly Unique?", Inches(0.5), Inches(0.4), Inches(12), Inches(0.8), sz=38, bold=True)

innovations = [
    ("🌐 Multilingual AI Feedback", "Beyond just text—AI reads reports aloud in Telugu/Hindi to bridge the literacy divide."),
    ("🌿 Hybrid Remedy Engine", "Bridges 5000 years of Ayurveda with modern Medicine; 80+ remedies for mild symptom management."),
    ("📍 1-Tap SOS Triage", "Automatically bypasses manual dialing; pressing 'CALL' instantly routes to 108 using native Web Protocols."),
    ("⚖️ Clinical Explainability", "Every risk percentage includes a transparent health explanation, not just a black-box number.")
]

for i, (t, b) in enumerate(innovations):
    xl = Inches(0.8); yl = Inches(1.8 + i*1.3)
    add_rect(s, xl, yl, Inches(11.7), Inches(1.1), fill=RGBColor(0x06,0x5f,0x46))
    add_rect(s, xl, yl, Inches(0.1), Inches(1.1), fill=C_HIGHLIGHT)
    add_txt(s, t, xl+Inches(0.3), yl+Inches(0.15), Inches(6), Inches(0.4), sz=18, bold=True, col=C_ACCENT)
    add_txt(s, b, xl+Inches(0.3), yl+Inches(0.6), Inches(11), Inches(0.4), sz=13, col=C_WHITE)

# ──────────────────────────────────────────────────────────
# SLIDE 6: SOCIAL IMPACT & SDG ALIGNMENT
# ──────────────────────────────────────────────────────────
s = blank_slide()
add_rect(s, 0, 0, W, H, fill=C_NAVY)
add_txt(s, "🌍 Scaling Toward a Healthier India", Inches(0.5), Inches(0.4), Inches(12), Inches(0.8), sz=36, bold=True)

metrics = [("SDG 3", "Good Health\n& Well-being"), ("SDG 10", "Reduced\nInequalities"), ("National", "Ayushman Bharat\nAlignment")]
for i, (n, l) in enumerate(metrics):
    xl = Inches(0.5 + i*4.3)
    add_rect(s, xl, Inches(1.5), Inches(4), Inches(2.2), fill=C_PRIMARY)
    add_txt(s, n, xl, Inches(1.8), Inches(4), Inches(0.8), sz=32, bold=True, col=C_HIGHLIGHT, align=PP_ALIGN.CENTER)
    add_txt(s, l, xl, Inches(2.7), Inches(4), Inches(0.6), sz=14, col=C_WHITE, align=PP_ALIGN.CENTER)

add_rect(s, Inches(0.5), Inches(4.2), Inches(12.3), Inches(2.2), fill=RGBColor(0x1e,0x29,0x3b))
add_txt(s, "Impact Roadmap", Inches(0.8), Inches(4.35), Inches(3), Inches(0.4), sz=18, bold=True, col=C_ACCENT)
impact = ["Phase 1: Full ASHA digitisation across 3 States (Hindi/Telugu belt).",
          "Phase 2: Integration with ABHA IDs (Ayushman Bharat Digital Mission).",
          "Phase 3: Native hospital connectivity via National Health Authority APIs."]
for i, text in enumerate(impact):
    add_txt(s, "• " + text, Inches(0.8), Inches(4.85 + i*0.45), Inches(11.5), Inches(0.4), sz=14)

# ──────────────────────────────────────────────────────────
# SLIDE 7: THANK YOU & LIVE DEMO
# ──────────────────────────────────────────────────────────
s = blank_slide()
add_rect(s, 0, 0, W, H, fill=RGBColor(0x06,0x1a,0x12))
add_img(s, IMG_COVER, Inches(6.5), 0, Inches(6.83), H)

add_txt(s, "🌱 JanRakshak AI", Inches(0.8), Inches(1.5), Inches(6), Inches(1.2), sz=52, bold=True)
add_txt(s, '"Predict Before It\'s Critical"', Inches(0.8), Inches(2.8), Inches(6), Inches(0.6), sz=22, col=C_ACCENT, italic=True)

add_rect(s, Inches(0.8), Inches(4.2), Inches(5.2), Inches(1.8), fill=C_PRIMARY)
add_txt(s, "🚀 Step into the Future", Inches(1.0), Inches(4.4), Inches(4.5), Inches(0.5), sz=16, bold=True, col=C_HIGHLIGHT)
add_txt(s, "Live Demo: http://localhost:5173", Inches(1.0), Inches(4.9), Inches(4.5), Inches(0.4), sz=13)
add_txt(s, "React 18 | Python Flask | Gemini AI", Inches(1.0), Inches(5.35), Inches(4.5), Inches(0.35), sz=11, col=C_GRAY)

add_txt(s, "Built with ❤️ for Rural India by Team JanRakshak", Inches(0.8), Inches(6.5), Inches(6), Inches(0.4), sz=14)

# ──────────────────────────────────────────────────────────
OUTPUT = r"D:\hackthon\JanRakshak_AI_Elite_Presentation.pptx"
prs.save(OUTPUT)
print(f"✅ Saved: {OUTPUT}")
