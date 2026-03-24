from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import pptx.oxml.ns as nsmap
from lxml import etree

# ── Colors ──────────────────────────────────────────────
DARK_GREEN   = RGBColor(0x1b, 0x43, 0x32)
MID_GREEN    = RGBColor(0x2d, 0x6a, 0x4f)
LIGHT_GREEN  = RGBColor(0x74, 0xc6, 0x9d)
MINT         = RGBColor(0xd1, 0xfa, 0xe5)
WHITE        = RGBColor(0xFF, 0xFF, 0xFF)
AMBER        = RGBColor(0xFF, 0x9F, 0x1C)
RED          = RGBColor(0xe6, 0x39, 0x46)
GRAY         = RGBColor(0x6b, 0x72, 0x80)
DARK         = RGBColor(0x1e, 0x29, 0x3b)

W = Inches(13.33)
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

def blank_slide(prs):
    blank_layout = prs.slide_layouts[6]
    return prs.slides.add_slide(blank_layout)

def bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def box(slide, l, t, w, h, fill_color=None, line_color=None, line_width=Pt(0)):
    shape = slide.shapes.add_shape(1, l, t, w, h)
    shape.line.width = line_width
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
    else:
        shape.line.fill.background()
    return shape

def txt(slide, text, l, t, w, h, size=24, bold=False, color=WHITE, align=PP_ALIGN.LEFT, wrap=True):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return tb

def heading(slide, text, top=Inches(0.35), size=36):
    txt(slide, text, Inches(0.5), top, Inches(12.3), Inches(0.7), size=size, bold=True, color=WHITE, align=PP_ALIGN.LEFT)

def subhead(slide, text, top=Inches(0.9), color=LIGHT_GREEN, size=16):
    txt(slide, text, Inches(0.5), top, Inches(12.3), Inches(0.4), size=size, bold=False, color=color, align=PP_ALIGN.LEFT)

def pill(slide, text, l, t, w, h, bg_color=MID_GREEN, text_color=WHITE, size=13):
    box(slide, l, t, w, h, fill_color=bg_color)
    txt(slide, text, l, t, w, h, size=size, bold=True, color=text_color, align=PP_ALIGN.CENTER)

def card(slide, title, body, l, t, w=Inches(3.8), h=Inches(1.6), title_color=LIGHT_GREEN, body_color=WHITE):
    box(slide, l, t, w, h, fill_color=MID_GREEN)
    txt(slide, title, l+Inches(0.15), t+Inches(0.1), w-Inches(0.3), Inches(0.45), size=14, bold=True, color=title_color)
    txt(slide, body,  l+Inches(0.15), t+Inches(0.5), w-Inches(0.3), h-Inches(0.55), size=11, color=body_color)

# ════════════════════════════════════════════════════════
# SLIDE 1 — COVER
# ════════════════════════════════════════════════════════
s = blank_slide(prs)
bg(s, DARK_GREEN)

# gradient bar top
box(s, 0, 0, W, Inches(0.15), fill_color=LIGHT_GREEN)
# gradient bar bottom
box(s, 0, H-Inches(0.15), W, Inches(0.15), fill_color=LIGHT_GREEN)

# Big title
txt(s, "🌱 JanRakshak AI", Inches(1), Inches(1.5), Inches(11.3), Inches(1.5),
    size=60, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# Tagline
txt(s, '"Predict Before It\'s Critical"', Inches(1), Inches(3.0), Inches(11.3), Inches(0.6),
    size=22, bold=False, color=LIGHT_GREEN, align=PP_ALIGN.CENTER)

# Sub
txt(s, "AI-Powered Rural Health Intelligence Platform", Inches(1), Inches(3.6), Inches(11.3), Inches(0.5),
    size=16, color=WHITE, align=PP_ALIGN.CENTER)

# Stats bar
stats = [("600M+","Rural Indians"), ("20+","Diseases"), ("3","Languages"), ("1-tap","Emergency Call")]
for i,(num,lbl) in enumerate(stats):
    xl = Inches(1.2 + i*2.8)
    box(s, xl, Inches(4.7), Inches(2.4), Inches(1.5), fill_color=MID_GREEN)
    txt(s, num, xl, Inches(4.75), Inches(2.4), Inches(0.7), size=26, bold=True, color=LIGHT_GREEN, align=PP_ALIGN.CENTER)
    txt(s, lbl, xl, Inches(5.4), Inches(2.4), Inches(0.5), size=12, color=WHITE, align=PP_ALIGN.CENTER)

# Badge
txt(s, "🏆  National Level Hackathon 2026  |  Healthcare & Social Impact",
    Inches(1), Inches(6.5), Inches(11.3), Inches(0.5), size=13, color=LIGHT_GREEN, align=PP_ALIGN.CENTER)

# ════════════════════════════════════════════════════════
# SLIDE 2 — THE PROBLEM
# ════════════════════════════════════════════════════════
s = blank_slide(prs)
bg(s, DARK)
box(s, 0, 0, Inches(0.18), H, fill_color=RED)

heading(s, "🔴  The Problem — India's Rural Health Crisis")
subhead(s, "600 million Indians lack basic primary healthcare access")

stats2 = [
    ("600M+", "Indians without\nprimary healthcare"),
    ("1.4M",  "Preventable deaths\nfrom late diagnosis/yr"),
    ("75%",   "PHCs understaffed /\nno rural doctors"),
    ("30+ km","Avg distance to\nnearest hospital"),
    ("~95%",  "ASHA workers with\nno digital tool"),
]
for i,(num,lbl) in enumerate(stats2):
    xl = Inches(0.4 + i*2.6)
    box(s, xl, Inches(1.5), Inches(2.4), Inches(2.3), fill_color=RGBColor(0x2d,0x1b,0x1b))
    box(s, xl, Inches(1.5), Inches(2.4), Inches(0.08), fill_color=RED)
    txt(s, num, xl, Inches(1.6), Inches(2.4), Inches(0.75), size=28, bold=True, color=RED, align=PP_ALIGN.CENTER)
    txt(s, lbl, xl, Inches(2.3), Inches(2.4), Inches(1.3), size=12, color=WHITE, align=PP_ALIGN.CENTER)

# Quote
box(s, Inches(0.5), Inches(4.2), Inches(12.3), Inches(1.4), fill_color=RGBColor(0x2a,0x2a,0x2a))
box(s, Inches(0.5), Inches(4.2), Inches(0.1), Inches(1.4), fill_color=RED)
txt(s, '"We knew she was sick, but we didn\'t know how sick. By the time she reached\nthe district hospital, it was too late." — ASHA Worker, Telangana (2024)',
    Inches(0.8), Inches(4.35), Inches(11.8), Inches(1.1), size=14, color=WHITE)

txt(s, "Paper registers → No risk scoring → No prescriptions → No GPS routing → Delayed treatment → Deaths",
    Inches(0.5), Inches(5.8), Inches(12.3), Inches(0.5), size=13, bold=True, color=AMBER, align=PP_ALIGN.CENTER)

# ════════════════════════════════════════════════════════
# SLIDE 3 — OUR SOLUTION
# ════════════════════════════════════════════════════════
s = blank_slide(prs)
bg(s, DARK_GREEN)

heading(s, "💡  Our Solution — JanRakshak AI")
subhead(s, "One app. Any smartphone. Instant AI-powered health screening for frontline workers.")

# Flow arrow
flow = ["Simple\nInputs", "AI Risk\nScore", "Medicine\nGuidance", "GPS Hospital\nRouting", "Emergency\nSupport"]
colors = [MID_GREEN, RGBColor(0x14,0x53,0x7a), RGBColor(0x14,0x53,0x7a), MID_GREEN, RED]
for i, (f, c) in enumerate(zip(flow, colors)):
    xl = Inches(0.35 + i*2.6)
    box(s, xl, Inches(1.5), Inches(2.2), Inches(1.1), fill_color=c)
    txt(s, f, xl, Inches(1.52), Inches(2.2), Inches(1.1), size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    if i < 4:
        txt(s, "→", xl+Inches(2.2), Inches(1.7), Inches(0.4), Inches(0.6), size=20, bold=True, color=LIGHT_GREEN, align=PP_ALIGN.CENTER)

# Feature grid
feats = [
    ("🧠 AI Risk Screening",     "Heart, Diabetes, Anemia scored\nwith %, color & AI explanation"),
    ("💊 Prescription Engine",   "Indian brand names auto-generated\n(Dolo, Glycomet, Dexorange)"),
    ("🔍 Medicine Lookup",       "20+ diseases, fuzzy keyword search\n— full dose/duration table"),
    ("🌿 Home Remedies",         "80+ traditional Indian remedies\nwith step-by-step instructions"),
    ("🗺️ GPS Hospital Map",      "Real location → Leaflet map\n+ 1.5km radius + hospitals"),
    ("🏥 1-Tap Emergency Call",  "108 · 102 · 112 + hospital numbers\ndial directly from the app"),
    ("🌐 3-Language Support",    "English · हिन्दी · తెలుగు\nfull UI translated instantly"),
    ("📊 ASHA Dashboard",        "Village analytics, recent records\nbulk screening mode"),
]
cols = 4
for i, (title, body) in enumerate(feats):
    col = i % cols
    row = i // cols
    xl = Inches(0.3 + col*3.26)
    yt = Inches(2.95 + row*1.75)
    card(s, title, body, xl, yt, w=Inches(3.1), h=Inches(1.6))

# ════════════════════════════════════════════════════════
# SLIDE 4 — TECHNICAL ARCHITECTURE
# ════════════════════════════════════════════════════════
s = blank_slide(prs)
bg(s, DARK)
box(s, 0, 0, Inches(0.18), H, fill_color=LIGHT_GREEN)

heading(s, "🏗️  Technical Architecture")

stack = [
    ("Frontend",  "React 18 + Vite",        "Fast, PWA-ready"),
    ("Styling",   "Vanilla CSS",             "Zero dependency"),
    ("Maps",      "Leaflet.js + OpenStreetMap","Free & offline-capable"),
    ("Backend",   "Python Flask + CORS",     "Lightweight API"),
    ("AI Engine", "Custom Scoring Algorithm","Transparent, explainable"),
    ("i18n",      "Custom translations.js",  "3-language, no lib"),
    ("Calling",   "tel: Web Protocol",       "Native phone dial"),
    ("Deploy",    "Vite + Flask server",      "Runs on 512MB RAM"),
]
for i, (layer, tech, reason) in enumerate(stack):
    col = i % 4
    row = i // 4
    xl = Inches(0.3 + col * 3.26)
    yt = Inches(1.3 + row * 1.5)
    box(s, xl, yt, Inches(3.1), Inches(1.35), fill_color=MID_GREEN)
    box(s, xl, yt, Inches(3.1), Inches(0.08), fill_color=LIGHT_GREEN)
    txt(s, layer, xl+Inches(0.12), yt+Inches(0.1), Inches(2.9), Inches(0.3), size=10, color=LIGHT_GREEN)
    txt(s, tech,  xl+Inches(0.12), yt+Inches(0.35), Inches(2.9), Inches(0.5), size=14, bold=True, color=WHITE)
    txt(s, reason,xl+Inches(0.12), yt+Inches(0.8), Inches(2.9), Inches(0.4), size=10, color=RGBColor(0xd1,0xd5,0xdb))

# API endpoints
box(s, Inches(0.3), Inches(4.55), Inches(12.7), Inches(1.8), fill_color=RGBColor(0x11,0x18,0x27))
txt(s, "API Endpoints", Inches(0.5), Inches(4.6), Inches(4), Inches(0.4), size=13, bold=True, color=LIGHT_GREEN)
txt(s, "POST  /api/predict       →  age, weight, BP, glucose, Hb, lat, lng  →  risk scores + color + explanation + prescriptions",
    Inches(0.5), Inches(5.0), Inches(12.3), Inches(0.4), size=11, color=WHITE)
txt(s, "POST  /api/disease-lookup  →  disease name (fuzzy)  →  medicines + brand names + dosage + home remedies",
    Inches(0.5), Inches(5.4), Inches(12.3), Inches(0.4), size=11, color=WHITE)

txt(s, "⚡ API response < 100ms  |  Bundle < 500KB  |  Runs on 512MB RAM / Raspberry Pi  |  No cloud dependency",
    Inches(0.3), Inches(6.5), Inches(12.7), Inches(0.5), size=12, bold=True, color=AMBER, align=PP_ALIGN.CENTER)

# ════════════════════════════════════════════════════════
# SLIDE 5 — INNOVATION
# ════════════════════════════════════════════════════════
s = blank_slide(prs)
bg(s, DARK_GREEN)

heading(s, "🚀  Innovation Highlights")

innovations = [
    ("🌐 Multilingual AI Output",      "Risk scores + explanations translated live into\nEnglish, Hindi, Telugu at the click of a button"),
    ("📍 GPS-Anchored Health Map",     "One-tap GPS → patient's real location → Leaflet\nmap + 1.5km radius + nearby hospital markers"),
    ("🔍 Fuzzy Medicine Search",       "Type 'loose motion' → matches Diarrhea\nType 'high BP' → matches Hypertension. Instant."),
    ("🌿 Traditional Remedy Engine",   "80+ Indian home remedies (Karela, Tulsi, Papaya\nLeaf, Giloy) with step-by-step preparation"),
    ("📞 1-Tap Emergency Calling",     "tel: protocol → pressing Call dials 108/102/112\nor hospital number directly. No app needed."),
    ("⚖️ Obesity-Adjusted AI Score",   "Weight >80kg adds obesity risk to diabetes score.\nClinically relevant, not just BMI-based."),
]
for i, (title, body) in enumerate(innovations):
    col = i % 3
    row = i // 3
    xl = Inches(0.3 + col * 4.35)
    yt = Inches(1.3 + row * 2.0)
    box(s, xl, yt, Inches(4.1), Inches(1.8), fill_color=MID_GREEN)
    box(s, xl, yt, Inches(0.08), Inches(1.8), fill_color=LIGHT_GREEN)
    txt(s, title, xl+Inches(0.2), yt+Inches(0.1), Inches(3.8), Inches(0.45), size=13, bold=True, color=LIGHT_GREEN)
    txt(s, body,  xl+Inches(0.2), yt+Inches(0.55), Inches(3.8), Inches(1.1), size=11, color=WHITE)

txt(s, "🏅  First-of-its-kind combination: AI Screening + Multilingual + GPS + Emergency Calling + Traditional Remedies — in a single rural health PWA",
    Inches(0.3), Inches(5.5), Inches(12.7), Inches(0.6), size=13, bold=True, color=AMBER, align=PP_ALIGN.CENTER)

# ════════════════════════════════════════════════════════
# SLIDE 6 — SOCIAL IMPACT
# ════════════════════════════════════════════════════════
s = blank_slide(prs)
bg(s, DARK)
box(s, 0, 0, Inches(0.18), H, fill_color=LIGHT_GREEN)

heading(s, "🌍  Social Impact & Government Alignment")

impacts = [
    ("1.2 Million", "ASHA workers\nempowered with AI"),
    ("600 Million", "Rural Indians who\nbenefit from early detection"),
    ("Immediate",   "Detection vs\n4–6 week wait"),
    ("100%",        "Screenings with\nprescription guidance"),
]
for i,(num,lbl) in enumerate(impacts):
    xl = Inches(0.35 + i*3.25)
    box(s, xl, Inches(1.2), Inches(3.0), Inches(1.5), fill_color=MID_GREEN)
    txt(s, num, xl, Inches(1.25), Inches(3.0), Inches(0.75), size=22, bold=True, color=LIGHT_GREEN, align=PP_ALIGN.CENTER)
    txt(s, lbl, xl, Inches(1.9), Inches(3.0), Inches(0.7), size=11, color=WHITE, align=PP_ALIGN.CENTER)

# Before/after table (text based)
box(s, Inches(0.3), Inches(3.0), Inches(12.7), Inches(0.35), fill_color=DARK_GREEN)
txt(s, "Metric", Inches(0.4), Inches(3.0), Inches(4), Inches(0.35), size=12, bold=True, color=WHITE)
txt(s, "Before JanRakshak AI", Inches(4.4), Inches(3.0), Inches(4), Inches(0.35), size=12, bold=True, color=WHITE)
txt(s, "After JanRakshak AI", Inches(8.5), Inches(3.0), Inches(4.2), Inches(0.35), size=12, bold=True, color=WHITE)

rows = [
    ("Time to detect high-risk patient",      "4–6 weeks (next camp)",  "✅ Immediate (same visit)"),
    ("Prescription guidance at screening",     "0%",                     "✅ 100% automated"),
    ("Emergency response",                     "30–45 min manual",       "✅ Instant 1-tap calling"),
    ("Medicine reference for ASHA workers",    "None",                   "✅ 20 diseases, 80 remedies"),
    ("Language options in health app",         "None",                   "✅ English, Hindi, Telugu"),
]
for i, (metric, before, after) in enumerate(rows):
    yt = Inches(3.4 + i * 0.55)
    fill = MID_GREEN if i % 2 == 0 else RGBColor(0x1a,0x35,0x28)
    box(s, Inches(0.3), yt, Inches(12.7), Inches(0.5), fill_color=fill)
    txt(s, metric, Inches(0.4), yt+Inches(0.05), Inches(4), Inches(0.4), size=11, color=WHITE)
    txt(s, before, Inches(4.4), yt+Inches(0.05), Inches(4), Inches(0.4), size=11, color=RGBColor(0xff,0x80,0x80))
    txt(s, after,  Inches(8.5), yt+Inches(0.05), Inches(4.2),Inches(0.4), size=11, color=LIGHT_GREEN)

txt(s, "SDG 3: Good Health  •  SDG 10: Reduced Inequalities  •  SDG 9: Innovation",
    Inches(0.3), Inches(6.5), Inches(7), Inches(0.5), size=12, color=LIGHT_GREEN)
txt(s, "Ayushman Bharat  •  NHM  •  e-Sanjeevani",
    Inches(7.3), Inches(6.5), Inches(5.5), Inches(0.5), size=12, color=AMBER, align=PP_ALIGN.RIGHT)

# ════════════════════════════════════════════════════════
# SLIDE 7 — FUTURE ROADMAP
# ════════════════════════════════════════════════════════
s = blank_slide(prs)
bg(s, DARK_GREEN)

heading(s, "🔭  Future Roadmap")

phases = [
    ("🟢 Phase 1\n0–3 months", MID_GREEN, [
        "Offline PWA (Service Workers)",
        "ABHA health ID integration",
        "Hindi/Telugu voice input",
        "Camera-based Hb estimation",
    ]),
    ("🔵 Phase 2\n3–6 months", RGBColor(0x14,0x53,0x7a), [
        "Real ML model (ICMR dataset)",
        "District-level heatmaps",
        "WhatsApp bot integration",
        "108 ambulance dispatch API",
    ]),
    ("🟡 Phase 3\n6–12 months", RGBColor(0x78,0x35,0x0f), [
        "Native Android app",
        "22 Indian language support",
        "NHA facility registry link",
        "Federated learning model",
    ]),
]
for i, (phase_title, color, items) in enumerate(phases):
    xl = Inches(0.35 + i * 4.35)
    box(s, xl, Inches(1.2), Inches(4.1), Inches(4.9), fill_color=color)
    txt(s, phase_title, xl+Inches(0.15), Inches(1.3), Inches(3.8), Inches(0.75), size=15, bold=True, color=WHITE)
    for j, item in enumerate(items):
        txt(s, f"→  {item}", xl+Inches(0.15), Inches(2.15)+Pt(j*35), Inches(3.8), Inches(0.5), size=12, color=WHITE)

txt(s, "💰  Business Model: Government NHM Contract  |  State SaaS License  |  NGO Partnership  |  Corporate CSR",
    Inches(0.3), Inches(6.5), Inches(12.7), Inches(0.5), size=13, bold=True, color=AMBER, align=PP_ALIGN.CENTER)

# ════════════════════════════════════════════════════════
# SLIDE 8 — WHY WE WILL WIN
# ════════════════════════════════════════════════════════
s = blank_slide(prs)
bg(s, DARK)

heading(s, "🎯  Why JanRakshak AI Will Win")

wins = [
    ("01", "Real Problem,\nReal Impact",   "Solves India's #1 healthcare gap. 600M users need this today — not a demo."),
    ("02", "Built for\nBharat",            "Multilingual, offline-capable, works on any Android. Designed for rural India."),
    ("03", "Complete\nProduct",            "AI + GPS + Prescriptions + Emergency + 20 diseases + 80 remedies. End-to-end."),
    ("04", "Government\nReady",            "Aligned with ASHA, NHM, Ayushman Bharat. Deploy-ready without rework."),
    ("05", "Technically\nDeep",            "Custom AI + Leaflet GPS + fuzzy search + tel: calling + real-time i18n."),
    ("06", "Live &\nDemonstrable",         "Judges can test the full app in 2 minutes. Not slides — a working product."),
]
for i, (num, title, body) in enumerate(wins):
    col = i % 3
    row = i // 3
    xl = Inches(0.35 + col * 4.35)
    yt = Inches(1.2 + row * 2.6)
    box(s, xl, yt, Inches(4.1), Inches(2.4), fill_color=MID_GREEN)
    box(s, xl, yt, Inches(4.1), Inches(0.1), fill_color=LIGHT_GREEN)
    txt(s, num,   xl+Inches(0.15), yt+Inches(0.15), Inches(3.8), Inches(0.65), size=32, bold=True, color=LIGHT_GREEN)
    txt(s, title, xl+Inches(0.15), yt+Inches(0.75), Inches(3.8), Inches(0.65), size=15, bold=True, color=WHITE)
    txt(s, body,  xl+Inches(0.15), yt+Inches(1.35), Inches(3.8), Inches(0.9), size=11, color=RGBColor(0xd1,0xd5,0xdb))

# ════════════════════════════════════════════════════════
# SLIDE 9 — CLOSING / THANK YOU
# ════════════════════════════════════════════════════════
s = blank_slide(prs)
bg(s, DARK_GREEN)
box(s, 0, 0, W, Inches(0.15), fill_color=LIGHT_GREEN)
box(s, 0, H-Inches(0.15), W, Inches(0.15), fill_color=LIGHT_GREEN)

txt(s, "🌱", Inches(5.9), Inches(0.8), Inches(1.5), Inches(1.0), size=52, align=PP_ALIGN.CENTER)
txt(s, "JanRakshak AI", Inches(1), Inches(1.7), Inches(11.3), Inches(1.0),
    size=52, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
txt(s, '"Swasth Bharat, Samridh Bharat"', Inches(1), Inches(2.8), Inches(11.3), Inches(0.6),
    size=20, color=LIGHT_GREEN, align=PP_ALIGN.CENTER)

txt(s, "A Healthy India is a Prosperous India", Inches(1), Inches(3.4), Inches(11.3), Inches(0.5),
    size=16, color=WHITE, align=PP_ALIGN.CENTER)

box(s, Inches(2), Inches(4.2), Inches(9.33), Inches(1.5), fill_color=MID_GREEN)
txt(s, "🚀  Live Demo: http://localhost:5176", Inches(2.2), Inches(4.3), Inches(9), Inches(0.5),
    size=16, bold=True, color=LIGHT_GREEN, align=PP_ALIGN.CENTER)
txt(s, "Frontend: React + Vite  |  Backend: Python Flask  |  AI: Custom Score Engine",
    Inches(2.2), Inches(4.8), Inches(9), Inches(0.5), size=13, color=WHITE, align=PP_ALIGN.CENTER)

txt(s, "Built with ❤️ for Rural India  •  © 2026 JanRakshak AI Team",
    Inches(1), Inches(6.3), Inches(11.3), Inches(0.5), size=13, color=LIGHT_GREEN, align=PP_ALIGN.CENTER)

# ── Save ─────────────────────────────────────────────────
output = r"D:\hackthon\JanRakshak_AI_Hackathon.pptx"
prs.save(output)
print(f"✅ Saved: {output}")
print(f"   Slides: {len(prs.slides)}")
